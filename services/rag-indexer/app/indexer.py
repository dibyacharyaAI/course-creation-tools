import logging
import asyncio
import os
from typing import List, Dict, Any
from pypdf import PdfReader
from pptx import Presentation

from shared.clients.embedding_client import GeminiEmbeddingClient
from shared.clients.vector_store_client import PGVectorClient
from .ocr_service import OCRService

logger = logging.getLogger(__name__)

class Indexer:
    def __init__(self, api_key: str, database_url: str, ocr_enabled: bool = False, deepseek_api_key: str = None):
        if not api_key:
            logger.warning("GEMINI_API_KEY is not set. Indexing will fail.")
        
        if api_key:
            self.embedding_client = GeminiEmbeddingClient(api_key=api_key)
        else:
            self.embedding_client = None
            logger.warning("GEMINI_API_KEY not set. Embedding client disabled.")
        
        if "sqlite" in database_url:
            logger.warning("Using SQLite. PGVectorClient disabled (Mock Mode).")
            class MockVectorStore:
                def add_documents(self, *args, **kwargs): pass
                def search(self, *args, **kwargs): return []
                def search_keyword(self, *args, **kwargs): 
                     # Mock keyword search
                     return [{"content": "Mock keyword result", "score": 1.0, "metadata": {}}]
            self.vector_store = MockVectorStore()
        else:
            self.vector_store = PGVectorClient(database_url=database_url)

        self.ocr_service = OCRService(api_key=deepseek_api_key, enabled=ocr_enabled)

    async def index_file(self, course_id: int, file_path: str, module_id: str = None, topic_id: str = None, extra_metadata: Dict[str, Any] = None):
        """
        Index a custom file (PDF, PPTX, TXT) for a course.
        Optional: Scope to specific module/topic.
        """
        try:
            text_content = ""
            filename = os.path.basename(file_path)
            ext = os.path.splitext(filename)[1].lower()

            logger.info(f"Processing file {filename} ({ext}) for course {course_id}")

            if ext == ".pdf":
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text_content += page.extract_text() + "\n"
            
            elif ext == ".pptx":
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
            
            elif ext == ".txt":
                with open(file_path, "r", encoding="utf-8") as f:
                    text_content = f.read()
            
            else:
                logger.warning(f"Unsupported file type: {ext}")
                return

            # Fallback to OCR if text is minimal (scanned PDF)
            if len(text_content.strip()) < 100:
                logger.info(f"Low text content detected ({len(text_content.strip())} chars). Attempting OCR fallback...")
                ocr_text = self.ocr_service.extract_text(file_path)
                if ocr_text:
                    text_content += "\n\n" + ocr_text
                    logger.info("OCR successfully added text content.")
                else:
                    logger.warning("OCR returned no text.")

            if not text_content.strip():
                logger.warning(f"No text extracted from {filename} after OCR attempt.")
                return

            # Simple chunking (e.g., 1000 chars)
            # In production, use a smarter chunker
            chunk_size = 1000
            chunks = [text_content[i:i+chunk_size] for i in range(0, len(text_content), chunk_size)]
            
            logger.info(f"Generated {len(chunks)} chunks from {filename}")

            # Create metadatas
            base_metadata = {
                "course_id": course_id,
                "source": filename,
                "type": "custom_upload"
            }
            if module_id:
                base_metadata["module_id"] = module_id
            if topic_id:
                base_metadata["topic_id"] = topic_id
            if extra_metadata:
                base_metadata.update(extra_metadata)
            
            metadatas = []
            for i in range(len(chunks)):
                m = base_metadata.copy()
                m["chunk_index"] = i
                metadatas.append(m)

            # Embed
            if not self.embedding_client:
                 raise ValueError("Embedding client not initialized (Missing API Key)")
            results = await self.embedding_client.embed_batch(chunks)
            embeddings = [r.embedding for r in results]

            # Store
            self.vector_store.add_documents(chunks, embeddings, metadatas)
            logger.info(f"✅ Successfully indexed custom file {filename} for course {course_id}")

        except Exception as e:
            logger.error(f"Failed to index file {file_path}: {e}")
            raise

    async def index_course_metadata(self, course_id: int, title: str, description: str):
        """
        Index basic course metadata (title, description).
        Useful for high-level relevance checks.
        """
        try:
            text = f"Course: {title}\nDescription: {description}"
            if not self.embedding_client:
                 raise ValueError("Embedding client not initialized (Missing API Key)")
            embedding = await self.embedding_client.embed(text, task_type="retrieval_document")
            
            self.vector_store.add_documents(
                texts=[text],
                embeddings=[embedding],
                metadatas=[{
                    "course_id": course_id,
                    "type": "course_metadata",
                    "title": title
                }]
            )
            logger.info(f"✅ Indexed metadata for course {course_id}")
        except Exception as e:
            logger.error(f"Failed to index course metadata for {course_id}: {e}")

    async def index_course_content(self, course_id: int, content: Dict[str, Any]):
        """
        Index full course content (modules, lessons).
        """
        try:
            chunks = []
            metadatas = []
            
            # Index Summary
            if "summary" in content:
                chunks.append(f"Course Summary: {content['summary']}")
                metadatas.append({
                    "course_id": course_id,
                    "type": "summary"
                })

            # Index Modules & Lessons
            if "modules" in content and isinstance(content["modules"], list):
                for module in content["modules"]:
                    module_id = module.get("code") or module.get("id") or "unknown"
                    module_title = module.get("title") or module.get("name") or ""
                    
                    # Index Module Description
                    module_text = f"Module {module_id}: {module_title}\n{module.get('description', '')}"
                    chunks.append(module_text)
                    metadatas.append({
                        "course_id": course_id,
                        "module_id": str(module_id),
                        "type": "module_overview",
                        "title": module_title
                    })
                    
                    # Index Lessons
                    if "lessons" in module and isinstance(module["lessons"], list):
                        for lesson in module["lessons"]:
                            lesson_code = lesson.get("code", "")
                            lesson_title = lesson.get("title", "")
                            lesson_body = lesson.get("body", "")
                            
                            # Chunk lesson body if too large (simple approach for now)
                            # Ideally split by paragraphs
                            lesson_text = f"Lesson {lesson_code}: {lesson_title}\n\n{lesson_body}"
                            
                            chunks.append(lesson_text)
                            metadatas.append({
                                "course_id": course_id,
                                "module_id": str(module_id),
                                "lesson_code": lesson_code,
                                "type": "lesson_content",
                                "title": lesson_title
                            })
            
            if not chunks:
                logger.warning(f"No content chunks found for course {course_id}")
                return

            logger.info(f"Generating embeddings for {len(chunks)} chunks...")
            if not self.embedding_client:
                 raise ValueError("Embedding client not initialized (Missing API Key)")
            results = await self.embedding_client.embed_batch(chunks)
            embeddings = [r.embedding for r in results]
            
            logger.info(f"Storing {len(chunks)} vectors...")
            self.vector_store.add_documents(chunks, embeddings, metadatas)
            logger.info(f"✅ Successfully indexed content for course {course_id}")
            
        except Exception as e:
            logger.error(f"Failed to index content for course {course_id}: {e}")
    async def hybrid_retrieve(self, course_id: int, query: str, k: int = 3) -> List[Dict[str, Any]]:
        """
        Retrieve using Reciprocal Rank Fusion (RRF) of Vector + Keyword search.
        """
        try:
            # 1. Get Top-K Vector Results
            if not self.embedding_client:
                 raise ValueError("Embedding client not initialized (Missing API Key)")

            query_vector = await self.embedding_client.embed(query, task_type="retrieval_query")
            vector_results = self.vector_store.search(
                query_vector=query_vector,
                top_k=k,
                filter={"course_id": str(course_id)}
            )

            # 2. Get Top-K Keyword Results
            keyword_results = self.vector_store.search_keyword(
                query_text=query,
                top_k=k,
                filter={"course_id": str(course_id)}
            )

            # 3. Apply RRF (Reciprocal Rank Fusion)
            # RRF score = 1 / (rank + k_rrf)
            k_rrf = 60
            fused_scores = {}
            doc_map = {}

            # Process Vector Results
            for rank, doc in enumerate(vector_results):
                doc_content = doc.get("content")
                doc_map[doc_content] = doc # Store full doc
                if doc_content not in fused_scores:
                    fused_scores[doc_content] = 0.0
                fused_scores[doc_content] += 1.0 / (rank + k_rrf)

            # Process Keyword Results
            for rank, doc in enumerate(keyword_results):
                doc_content = doc.get("content")
                if doc_content not in doc_map:
                    doc_map[doc_content] = doc
                if doc_content not in fused_scores:
                    fused_scores[doc_content] = 0.0
                fused_scores[doc_content] += 1.0 / (rank + k_rrf)

            # Sort by RRF Score
            sorted_docs = sorted(fused_scores.items(), key=lambda item: item[1], reverse=True)
            
            # Return Top-K
            final_results = []
            for content, score in sorted_docs[:k]:
                doc = doc_map[content]
                doc["rrf_score"] = score
                final_results.append(doc)
            
            return final_results

        except Exception as e:
            logger.error(f"Hybrid retrieval failed for course {course_id}: {e}")
            return []

    async def retrieve(self, course_id: int, query: str, k: int = 3) -> List[Dict[str, Any]]:
        """
        Retrieve relevant snippets for a given query within a course scope.
        """
        try:
            # 1. Embed Query
            query_vector = await self.embedding_client.embed(query, task_type="retrieval_query")
            
            # 2. Search Vector Store
            results = self.vector_store.search(
                query_vector=query_vector,
                top_k=k,
                filter={"course_id": str(course_id)}
            )
            return results
            
        except Exception as e:
            logger.error(f"Retrieval failed for course {course_id} query '{query}': {e}")
            return []
