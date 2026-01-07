import os
import logging
from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

class PPTXGenerator:
    def __init__(self, output_dir: str = "/app/exports"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def generate_slides(self, course_data: dict) -> str:
        """
        Generate PPTX slides from course data.
        Returns the filename of the generated PPTX.
        """
        try:
            prs = Presentation()
            
            # Title Slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = course_data.get('title', 'Untitled Course')
            subtitle.text = f"{course_data.get('programme', '')} - {course_data.get('semester', '')}\n{course_data.get('course_code', '')}"
            
            # Content Slides
            modules = course_data.get('content', {}).get('modules', [])
            for module in modules:
                # Module Title Slide
                layout = prs.slide_layouts[2] # Section Header
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = f"{module.get('code', '')}: {module.get('title', '')}"
                
                for lesson in module.get('lessons', []):
                    slide_outline = lesson.get('slide_outline', {}).get('slides', [])
                    
                    for slide_data in slide_outline:
                        bullet_layout = prs.slide_layouts[1] # Title and Content
                        slide = prs.slides.add_slide(bullet_layout)
                        
                        # Set Title
                        slide.shapes.title.text = slide_data.get('title', lesson.get('title', ''))
                        
                        # Add Bullets
                        tf = slide.placeholders[1].text_frame
                        bullets = slide_data.get('bullets', [])
                        
                        if bullets:
                            tf.text = bullets[0]
                            for bullet in bullets[1:]:
                                p = tf.add_paragraph()
                                p.text = bullet
                        
                        # Add Footer with CO mapping
                        if 'co_mapping' in lesson:
                            footer_txt = f"Mapped to: {', '.join(lesson['co_mapping'])}"
                            # Add text box for footer
                            left = Inches(0.5)
                            top = Inches(7)
                            width = Inches(9)
                            height = Inches(0.5)
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.text_frame
                            p = tf.add_paragraph()
                            p.text = footer_txt
                            p.font.size = Pt(10)
                            p.font.color.rgb = None # Default color
            
            filename = f"course_{course_data['id']}_slides.pptx"
            output_path = os.path.join(self.output_dir, filename)
            
            prs.save(output_path)
            logger.info(f"Generated PPTX: {output_path}")
            
            return filename
            
        except Exception as e:
            logger.error(f"Failed to generate PPTX: {e}")
            raise e
